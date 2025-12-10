import streamlit as st
import streamlit.components.v1 as components
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE 
from io import BytesIO
import docx
from docx.document import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
import fitz  # PyMuPDF
import re
import pandas as pd
from PIL import Image
import pytesseract

# --- è¨­å®šç¶²é æ¨™é¡Œ ---
st.set_page_config(page_title="PPT é‡çµ„ç”Ÿæˆå™¨ (V17 å¯¬é¬†æ¯”å°ç‰ˆ)", page_icon="ğŸ“‘", layout="wide")
st.title("ğŸ“‘ PPT é‡çµ„ç”Ÿæˆå™¨ (V17 å¯¬é¬†æ¯”å°ç‰ˆ)")
st.caption("æ›´æ–°ï¼šV17 åŠ å…¥ã€Œæ ¸å¿ƒæ•¸å­—æ¯”å°ã€æ©Ÿåˆ¶ã€‚ç•¶ PDF æª”ååŒ…å«é¡å¤–çš„é›¶ (å¦‚ us000123...) æ™‚ï¼Œç¨‹å¼èƒ½é€éè­˜åˆ¥é—œéµæ•¸å­—ä¸² (å¦‚ 123) æˆåŠŸæ‰¾åˆ°å°æ‡‰æª”æ¡ˆã€‚")

# === NBLM æç¤ºè©å€å¡Š ===
nblm_prompt = """æ ¹æ“šä¸Šå‚³çš„æ‰€æœ‰ä¾†æºï¼Œåˆ†é–‹æ•´ç†å‡ºä»¥ä¸‹é‡é»(ä¸è¦è¡¨æ ¼)ï¼š

1. æ¡ˆè™Ÿ / æ—¥æœŸ / å…¬å¸ï¼š *(æ¡ˆè™Ÿä¾æ“š"å…¬é–‹è™Ÿ"ã€æ—¥æœŸä¾æ“š"å„ªå…ˆæ¬Šæ—¥"ã€å…¬å¸ä¾æ“š"ç”³è«‹äºº")
2. è§£æ±ºå•é¡Œï¼š
3. ç™¼æ˜ç²¾ç¥ï¼š*(ä¸è¦æœ‰å…¬å¼)
4. ä¸€å¥é‡é»ï¼š *(ç”¨ä¾†æè¿°ç™¼æ˜ç‰¹å¾µé‡é»ï¼Œ20å­—)
5. ä»£è¡¨åœ–ï¼š*(æ ¹æ“šç™¼æ˜ç²¾ç¥å»ºè­°3å¼µæœ€å¯ä»¥èªªæ˜ç™¼æ˜ç²¾ç¥çš„åœ–ç‰‡ï¼Œç¯„ä¾‹:FIG.3)
6. ç¨ç«‹é …claimï¼š *(åˆ†çµ„ä¸”åˆ†è¡Œæ¢åˆ—å¼+å°æ‡‰çš„ä»£è¡¨åœ–ï¼Œclaimè¦(1)æœ‰ä½éšç¸®æ’ (2)claimçš„å…ƒä»¶è¦æœ‰æ¨™è™Ÿ (3)å°æ‡‰çš„claimè™Ÿç¢¼)"""

st.info("ğŸ’¡ **NBLM ä½¿ç”¨æç¤ºè©** (é»æ“Šä¸‹æ–¹ç¶ è‰²æŒ‰éˆ•ä¸€éµè¤‡è£½)")

components.html(
    f"""
    <html>
    <head><meta charset="utf-8"></head>
    <body style="font-family: sans-serif; margin: 0; padding: 0;">
        <div style="display: flex; flex-direction: column; align-items: flex-start;">
            <textarea id="copyTarget" style="opacity: 0; position: absolute; z-index: -1;">{nblm_prompt}</textarea>
            <div style="background-color: #f0f2f6; padding: 15px; border-radius: 10px; white-space: pre-wrap; font-size: 14px; color: #31333F; border: 1px solid #d6d6d6; width: 95%; margin-bottom: 10px;">{nblm_prompt}</div>
            <button onclick="copyFunction()" style="background-color: #00CC66; color: white; border: none; padding: 12px 24px; font-size: 16px; font-weight: bold; border-radius: 8px; cursor: pointer; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">ğŸ“‹ é»æˆ‘ä¸€éµè¤‡è£½æç¤ºè©</button>
            <span id="statusParams" style="color: #00CC66; font-weight: bold; margin-left: 10px; opacity: 0; transition: opacity 0.5s;">âœ… è¤‡è£½æˆåŠŸï¼</span>
        </div>
        <script>
        function copyFunction() {{
            var copyText = document.getElementById("copyTarget");
            copyText.select();
            navigator.clipboard.writeText(copyText.value).then(function() {{
                var status = document.getElementById("statusParams");
                status.style.opacity = '1';
                setTimeout(function(){{ status.style.opacity = '0'; }}, 2000);
            }});
        }}
        </script>
    </body>
    </html>
    """,
    height=360
)
st.divider()

# --- åˆå§‹åŒ– Session State ---
if 'slides_data' not in st.session_state:
    st.session_state['slides_data'] = []
if 'status_report' not in st.session_state:
    st.session_state['status_report'] = []
if 'debug_logs_map' not in st.session_state:
    st.session_state['debug_logs_map'] = {}

# --- è¼”åŠ©å‡½æ•¸ï¼šéæ­· Word ---
def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    else:
        raise ValueError("åªæ”¯æ´è®€å–æ•´ä»½ Document")
    for child in parent_elm.iterchildren():
        if child.tag.endswith('p'):
            yield Paragraph(child, parent)
        elif child.tag.endswith('tbl'):
            yield Table(child, parent)

# --- æ ¸å¿ƒå‡½æ•¸ï¼šV13 å‹•æ…‹å¹³è¡¡ç‰ˆé‚è¼¯ ---
def extract_images_from_pdf_v13(pdf_stream, target_fig_text, case_key, debug=False, log_prefix=""):
    if not target_fig_text:
        return [], f"{log_prefix}æœªæŒ‡å®šåœ–è™Ÿ"
    
    try:
        if hasattr(pdf_stream, 'seek'):
            pdf_stream.seek(0)
            
        doc = fitz.open(stream=pdf_stream, filetype="pdf")
        
        matches = re.findall(r'(?:FIG\.?|Figure|å›¾|åœ–)[\s\.]*([0-9]+[A-Za-z]*)', target_fig_text, re.IGNORECASE)
        if not matches:
            first_line = target_fig_text.split('\n')[0].strip().upper()
            fallback = re.search(r'([0-9]+[A-Z]*)', first_line)
            if fallback: matches = [fallback.group(1)]

        if not matches:
            return [], f"{log_prefix}ç„¡æ³•è­˜åˆ¥ä»»ä½•åœ–è™Ÿ"

        target_numbers = sorted(list(set([m.upper() for m in matches])))
        
        # V13 åƒæ•¸
        PAGE_TEXT_THRESHOLD_OCR = 800  
        PAGE_TEXT_THRESHOLD_RAW = 600 
        LONG_SENTENCE_LIMIT = 80 
        MAX_LONG_SENTENCES = 3
        LINE_LENGTH_LIMIT = 30

        page_blacklist_headers = [
            "BRIEF DESCRIPTION", "DETAILED DESCRIPTION", "å…·ä½“å®æ–½æ–¹å¼", "å¯¦æ–½æ–¹å¼", 
            "WHAT IS CLAIMED", "æƒåˆ©è¦æ±‚", "ç”³è«‹å°ˆåˆ©ç¯„åœ", "åœ–å¼ç°¡å–®èªªæ˜", "ã€åœ–å¼ç°¡å–®èªªæ˜ã€‘",
            "ABSTRACT", "æ‘˜è¦", "BACKGROUND", "èƒŒæ™¯æŠ€è¡“",
            "ç¬¦å·è¯´æ˜", "ç¬¦è™Ÿèªªæ˜"
        ]

        SENTENCE_STOPWORDS = ["ç‚º", "ä¿‚", "æ‰€ç¤º", "é—œæ–¼", "åƒç…§", "åƒè€ƒ", "EXAMPLE", "EMBODIMENT", "SHOWS", "REFER"]

        found_page_indices = set()
        debug_logs = [] 
        debug_logs.append(f"{log_prefix}ğŸ¯ ç›®æ¨™: {target_numbers}")

        for target_number in target_numbers:
            search_tokens = [
                f"FIG{target_number}", f"FIGURE{target_number}",
                f"å›¾{target_number}", f"åœ–{target_number}"
            ]
            
            found_this_fig = False

            for i, page in enumerate(doc):
                blocks = page.get_text("blocks")
                page_text_all = "".join([b[4] for b in blocks]).upper()
                clean_page_text_all = re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '', page_text_all)
                page_text_len = len(clean_page_text_all)

                is_blacklist_page = False
                for header in page_blacklist_headers:
                    if header in page_text_all:
                        is_blacklist_page = True
                        if debug and i < 15: debug_logs.append(f"{log_prefix}ğŸš« Skip P{i+1} (Header: {header})")
                        break
                if is_blacklist_page: continue

                long_sentence_count = 0
                for b in blocks:
                    if len(re.sub(r'\s+', '', b[4])) > LONG_SENTENCE_LIMIT:
                        long_sentence_count += 1
                
                if long_sentence_count > MAX_LONG_SENTENCES:
                    if debug and i < 15: debug_logs.append(f"{log_prefix}ğŸš« Skip P{i+1} (Raw: Long sentences)")
                    continue

                if page_text_len > PAGE_TEXT_THRESHOLD_RAW:
                    if debug and i < 15: debug_logs.append(f"{log_prefix}ğŸš« Skip P{i+1} (Raw Heavy: {page_text_len})")
                    continue

                match_found_strategy_1 = False
                for b in blocks:
                    block_text = b[4].strip()
                    clean_block_text = re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '', block_text).upper()
                    
                    for token in search_tokens:
                        if token in clean_block_text:
                            if len(clean_block_text) > LINE_LENGTH_LIMIT: continue
                            
                            is_sentence = False
                            for stopword in SENTENCE_STOPWORDS:
                                if stopword in clean_block_text:
                                    is_sentence = True
                                    break
                            if is_sentence: continue 

                            idx = clean_block_text.find(token)
                            is_exact_match = True
                            if idx != -1:
                                after_idx = idx + len(token)
                                if after_idx < len(clean_block_text) and clean_block_text[after_idx].isdigit():
                                    is_exact_match = False
                            
                            if is_exact_match:
                                found_page_indices.add(i)
                                found_this_fig = True
                                match_found_strategy_1 = True
                                if debug: debug_logs.append(f"{log_prefix}âœ… Found {token} (Text Layer) on P{i+1}")
                                break
                    if match_found_strategy_1: break
                
                if match_found_strategy_1: 
                    if found_this_fig: break
                    continue

                if page_text_len < PAGE_TEXT_THRESHOLD_RAW:
                    for token in search_tokens:
                        if token in clean_page_text_all:
                            idx = clean_page_text_all.find(token)
                            is_exact_match = True
                            if idx != -1:
                                after_idx = idx + len(token)
                                if after_idx < len(clean_page_text_all) and clean_page_text_all[after_idx].isdigit():
                                    is_exact_match = False
                            if is_exact_match:
                                found_page_indices.add(i)
                                found_this_fig = True
                                match_found_strategy_1 = True
                                if debug: debug_logs.append(f"{log_prefix}âœ… Found {token} (Full Page) on P{i+1}")
                                break
                
                if match_found_strategy_1:
                    if found_this_fig: break
                    continue

                if page_text_len < 200: 
                    try:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_data = pix.tobytes("png")
                        pil_image = Image.open(BytesIO(img_data))
                        
                        ocr_text = pytesseract.image_to_string(pil_image, lang='eng+chi_tra', config='--psm 11')
                        ocr_text_clean = re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '', ocr_text).upper()
                        ocr_len = len(ocr_text_clean)

                        if debug and i < 15: debug_logs.append(f"{log_prefix}ğŸ‘ï¸ OCR P{i+1} Len: {ocr_len}")

                        if ocr_len > PAGE_TEXT_THRESHOLD_OCR:
                            if debug: debug_logs.append(f"{log_prefix}   -> Skip P{i+1} (OCR Heavy)")
                            continue
                        
                        ocr_lines = ocr_text.split('\n')
                        long_sentence_count_ocr = 0
                        for line in ocr_lines:
                             clean_line_len = len(re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '', line))
                             if clean_line_len > LONG_SENTENCE_LIMIT:
                                 long_sentence_count_ocr += 1
                        
                        if long_sentence_count_ocr > MAX_LONG_SENTENCES:
                             if debug: debug_logs.append(f"{log_prefix}   -> Skip P{i+1} (OCR Long Sentences)")
                             continue

                        for line in ocr_lines:
                            clean_line = re.sub(r'[^a-zA-Z0-9\u4e00-\u9fa5]', '', line).upper()
                            
                            for token in search_tokens:
                                if token in clean_line:
                                    if len(clean_line) > LINE_LENGTH_LIMIT: continue
                                    
                                    is_sentence_ocr = False
                                    for stopword in SENTENCE_STOPWORDS:
                                        if stopword in clean_line:
                                            is_sentence_ocr = True
                                            break
                                    if is_sentence_ocr: continue

                                    found_page_indices.add(i)
                                    found_this_fig = True
                                    if debug: debug_logs.append(f"{log_prefix}âœ… Found {token} (OCR) on P{i+1}")
                                    break
                            if found_this_fig: break

                    except Exception as ocr_e:
                        if debug: debug_logs.append(f"{log_prefix}âš ï¸ OCR Error on P{i+1}: {ocr_e}")

                if found_this_fig: break
        
        if debug:
            if case_key not in st.session_state['debug_logs_map']:
                st.session_state['debug_logs_map'][case_key] = ""
            st.session_state['debug_logs_map'][case_key] += "\n".join(debug_logs) + "\n\n"

        if not found_page_indices:
            return [], f"{log_prefix}æ‰¾ä¸åˆ°åœ–è™Ÿ: {', '.join(target_numbers)}"

        output_images = []
        for page_idx in sorted(list(found_page_indices)):
            page = doc[page_idx]
            mat = fitz.Matrix(4, 4) 
            pix = page.get_pixmap(matrix=mat)
            output_images.append(pix.tobytes("png"))

        return output_images, f"æˆåŠŸ ({len(output_images)}å¼µ)"

    except Exception as e:
        return [], f"{log_prefix}PDF è§£æéŒ¯èª¤: {str(e)}"

# --- å‡½æ•¸ï¼šæå–å°ˆåˆ©è™Ÿ (V16 ä¿®æ­£ï¼šå»é™¤é€—è™Ÿ) ---
def extract_patent_number_from_text(text):
    if "ï¼š" in text: text = text.replace("ï¼š", ":")
    if ":" in text:
        content = text.split(":", 1)[1]
    else:
        content = text

    clean_text = content.replace(" ", "").replace(",", "").strip().upper()
    match = re.search(r'([A-Z]{2,4}\d{4,}[A-Z0-9]*)', clean_text)
    if match: return match.group(1)
    
    return ""

def extract_header_info_detail(raw_text):
    number = "(æœªæ‰¾åˆ°)"
    date = "(æœªæ‰¾åˆ°)"
    company = "(æœªæ‰¾åˆ°)"
    
    extracted_no = extract_patent_number_from_text(raw_text)
    if extracted_no: number = extracted_no
    else:
        match_no = re.search(r'(?:å…¬é–‹è™Ÿ|æ¡ˆè™Ÿ)[:ï¼š\s]*([^\n]+)', raw_text)
        if match_no:
            raw_no = match_no.group(1)
            raw_no = re.split(r'\s+(?:æ—¥æœŸ|å…¬å¸|ç”³è«‹äºº)[:ï¼š]', raw_no)[0]
            number = raw_no.strip()

    match_date = re.search(r'(?:æ—¥æœŸ)[:ï¼š\s]*(\d{4}[./-]\d{1,2}[./-]\d{1,2})', raw_text)
    if match_date: date = match_date.group(1).strip()
    else:
        match_date_backup = re.search(r'(\d{4}[./-]\d{1,2}[./-]\d{1,2})', raw_text)
        if match_date_backup: date = match_date_backup.group(1).strip()

    matches = re.findall(r'(?:å…¬å¸|ç”³è«‹äºº)[:ï¼š\s]*(.*?)(?=\s+(?:å…¬é–‹è™Ÿ|æ¡ˆè™Ÿ|æ—¥æœŸ)[:ï¼š]|$)', raw_text)
    if matches:
        for candidate in reversed(matches):
            clean_cand = candidate.strip()
            if len(clean_cand) > 1 and "å…¬é–‹è™Ÿ" not in clean_cand:
                company = clean_cand
                break

    return number, date, company

def extract_date_for_sort(text):
    match = re.search(r'(\d{4})[./-](\d{1,2})[./-](\d{1,2})', text)
    if match: return f"{match.group(1)}{match.group(2).zfill(2)}{match.group(3).zfill(2)}"
    return "99999999"

def extract_company_for_sort(text):
    _, _, comp = extract_header_info_detail(text)
    if comp != "(æœªæ‰¾åˆ°)": return comp
    return "ZZZ"

def normalize_string(s):
    if not s: return ""
    return re.sub(r'[^A-Z0-9]', '', s.upper())

def parse_word_file(uploaded_docx):
    try:
        doc = docx.Document(uploaded_docx)
        cases = []
        current_case = {
            "case_info": "", "problem": "", "spirit": "", "key_point": "", "rep_fig_text": "", "claim_text": "",
            "image_list": [], "claim_image_list": [], "image_name": "WordåŒ¯å…¥", "raw_case_no": "",
            "clean_number": "", "clean_date": "", "clean_company": "", 
            "sort_date": "99999999", "sort_company": "ZZZ",
            "source_file": uploaded_docx.name, "missing_fields": []
        }
        current_field = None 
        
        all_lines = []
        for block in iter_block_items(doc):
            if isinstance(block, Paragraph):
                if block.text.strip(): all_lines.append(block.text.strip())
            elif isinstance(block, Table):
                for row in block.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            if p.text.strip(): all_lines.append(p.text.strip())
        
        for text in all_lines:
            if "æ¡ˆè™Ÿ" in text or "ç´¢è™Ÿ" in text:
                if current_case["case_info"] and current_field != "case_info_block":
                    nb, dt, cp = extract_header_info_detail(current_case["case_info"])
                    current_case["clean_number"] = nb
                    current_case["clean_date"] = dt
                    current_case["clean_company"] = cp
                    if not current_case["problem"]: current_case["missing_fields"].append("è§£æ±ºå•é¡Œ")
                    cases.append(current_case)
                    current_case = {
                        "case_info": "", "problem": "", "spirit": "", "key_point": "", "rep_fig_text": "", "claim_text": "",
                        "image_list": [], "claim_image_list": [], "image_name": "WordåŒ¯å…¥", "raw_case_no": "",
                        "clean_number": "", "clean_date": "", "clean_company": "",
                        "sort_date": "99999999", "sort_company": "ZZZ",
                        "source_file": uploaded_docx.name, "missing_fields": []
                    }
                current_field = "case_info_block"
                current_case["case_info"] = text
                
                nb, dt, cp = extract_header_info_detail(text)
                if dt != "(æœªæ‰¾åˆ°)": current_case["sort_date"] = dt.replace(".", "").replace("/", "").replace("-", "")
                if cp != "(æœªæ‰¾åˆ°)": current_case["sort_company"] = cp
                if nb != "(æœªæ‰¾åˆ°)": current_case["raw_case_no"] = nb
                continue

            if "è§£æ±ºå•é¡Œ" in text:
                current_field = "problem"
                current_case["problem"] = re.sub(r'^[0-9.ï¼]*\s*è§£æ±ºå•é¡Œ[:ï¼š]?\s*', '', text)
                continue
            elif "ç™¼æ˜ç²¾ç¥" in text:
                current_field = "spirit"
                current_case["spirit"] = re.sub(r'^[0-9.ï¼]*\s*ç™¼æ˜ç²¾ç¥[:ï¼š]?\s*', '', text)
                continue
            elif "é‡é»" in text:
                current_field = "key_point"
                current_case["key_point"] = re.sub(r'^[0-9.ï¼]*\s*(ä¸€å¥)?é‡é»[:ï¼š]?\s*', '', text)
                continue
            elif "ä»£è¡¨åœ–" in text:
                current_field = "rep_fig"
                current_case["rep_fig_text"] = re.sub(r'^[0-9.ï¼]*\s*ä»£è¡¨åœ–[:ï¼š]?\s*', '', text).strip()
                continue
            elif "ç¨ç«‹é …" in text or ("claim" in text.lower() and "6" in text):
                current_field = "claim"
                content = re.sub(r'^[0-9.ï¼]*\s*(ç¨ç«‹é …)?(claim)?[:ï¼š]?\s*', '', text, flags=re.IGNORECASE).strip()
                current_case["claim_text"] = content
                continue

            if current_field == "case_info_block":
                current_case["case_info"] += "\n" + text
                nb, dt, cp = extract_header_info_detail(current_case["case_info"])
                if dt != "(æœªæ‰¾åˆ°)": current_case["sort_date"] = dt.replace(".", "").replace("/", "").replace("-", "")
                if cp != "(æœªæ‰¾åˆ°)": current_case["sort_company"] = cp
                if nb != "(æœªæ‰¾åˆ°)": current_case["raw_case_no"] = nb
            elif current_field == "rep_fig":
                current_case["rep_fig_text"] += "\n" + text
            elif current_field == "problem":
                current_case["problem"] += "\n" + text
            elif current_field == "spirit":
                current_case["spirit"] += "\n" + text
            elif current_field == "key_point":
                current_case["key_point"] += "\n" + text
            elif current_field == "claim": 
                current_case["claim_text"] += "\n" + text

        if current_case["case_info"]:
            nb, dt, cp = extract_header_info_detail(current_case["case_info"])
            current_case["clean_number"] = nb
            current_case["clean_date"] = dt
            current_case["clean_company"] = cp
            if not current_case["problem"]: current_case["missing_fields"].append("è§£æ±ºå•é¡Œ")
            cases.append(current_case)
        return cases
    except Exception as e:
        st.error(f"è§£æ Word éŒ¯èª¤ ({uploaded_docx.name}): {e}")
        return []

def split_claims_text(full_text):
    if not full_text: return []
    lines = full_text.split('\n')
    claims = []
    current_chunk = []
    header_pattern = re.compile(r'(\(Claim\s*\d+\)|^\s*(Claim|ç¨ç«‹é …)\s*\d+|^\s*\d+\.\s)', re.IGNORECASE)
    for line in lines:
        if header_pattern.search(line):
            if current_chunk:
                if "".join(current_chunk).strip(): claims.append(current_chunk)
            current_chunk = [line]
        else:
            current_chunk.append(line)
    if current_chunk and "".join(current_chunk).strip(): claims.append(current_chunk)
    return claims

def parse_fig_number_from_claim(claim_text):
    if not claim_text: return None
    matches = re.findall(r'(?:FIG\.?|Figure|å›¾|åœ–)[\s\.]*([0-9]+[A-Za-z]*)', claim_text, re.IGNORECASE)
    if matches:
        return "FIG. " + ", FIG. ".join(sorted(list(set(matches))))
    return None

# --- å´é‚Šæ¬„ ---
with st.sidebar:
    st.header("1. åŒ¯å…¥è³‡æ–™")
    word_files = st.file_uploader("Word æª”æ¡ˆ (å¯å¤šé¸)", type=['docx'], accept_multiple_files=True)
    pdf_files = st.file_uploader("PDF æª”æ¡ˆ (å¯å¤šé¸)", type=['pdf'], accept_multiple_files=True)
    st.divider()
    st.header("2. è¼¸å‡ºè¨­å®š")
    add_claim_slide = st.checkbox("âœ… æ˜¯å¦ç”¢ç”Ÿ Claim åˆ†é ", value=False, help="å‹¾é¸å¾Œï¼Œç¨‹å¼æœƒè‡ªå‹•è­˜åˆ¥ç¨ç«‹é …æ•¸é‡ï¼Œä¸¦ç‚ºæ¯ä¸€çµ„ç¨ç«‹é …ç”¢ç”Ÿä¸€é ")
    
    st.divider()
    st.header("3. é€²éšé™¤éŒ¯")
    debug_mode = st.checkbox("ğŸ é–‹å•ŸåµéŒ¯æ¨¡å¼ (Debug)", value=False, help="å‹¾é¸å¾Œï¼Œæœƒé¡¯ç¤ºè©³ç´°çš„è­˜åˆ¥æ—¥èªŒï¼ŒåŒ…å« OCR çš„è¾¨è­˜çµæœã€‚")
    
    if debug_mode and st.session_state['debug_logs_map']:
        st.caption("ğŸ“œ æ­·å² Debug ç´€éŒ„ (é»æ“Šå±•é–‹)")
        for key, log in st.session_state['debug_logs_map'].items():
            with st.expander(f"Case: {key}"):
                st.text(log)

    st.divider()
    run_btn = st.button("ğŸ”„ é–‹å§‹æ™ºèƒ½æ•´åˆ", type="primary")

    if run_btn:
        st.session_state['debug_logs_map'] = {}
        
        if not word_files:
            st.warning("âš ï¸ è«‹å…ˆä¸Šå‚³ Word æª”æ¡ˆï¼")
            st.stop()

        all_cases = []
        status_report_list = []
        for wf in word_files: all_cases.extend(parse_word_file(wf))
        
        pdf_file_map = {}
        if pdf_files:
            for pf in pdf_files:
                pdf_file_map[pf.name] = pf.read()

        match_count = 0
        current_ppt_page = 1 
        with st.spinner("è™•ç†ä¸­... (è‹¥å•Ÿå‹• OCR å¯èƒ½éœ€è¦è¼ƒé•·æ™‚é–“ï¼Œè«‹è€å¿ƒç­‰å€™)"):
            all_cases.sort(key=lambda x: (x["sort_company"].upper(), x["sort_date"]))
            for case in all_cases:
                case_key = case["raw_case_no"]
                target_fig = case["rep_fig_text"]
                claim_text_content = case["claim_text"]
                
                pages_this_case = 1 
                if add_claim_slide:
                    c_groups = split_claims_text(claim_text_content)
                    if not c_groups and claim_text_content.strip(): pages_this_case += 1
                    else: pages_this_case += len(c_groups)
                
                start_page = current_ppt_page
                end_page = current_ppt_page + pages_this_case - 1
                page_str = f"P{start_page}" if start_page == end_page else f"P{start_page}-P{end_page}"
                current_ppt_page += pages_this_case

                status = {
                    "ä¾†æº": case["source_file"], 
                    "æ¡ˆè™Ÿ(å…¬é–‹è™Ÿ)": case["clean_number"],
                    "å…¬å¸": case["clean_company"],
                    "æ—¥æœŸ(å„ªå…ˆæ¬Šæ—¥)": case["clean_date"],
                    "å°æ‡‰PPTçš„é ç¢¼": page_str,
                    "ç‹€æ…‹": "æœªè™•ç†", "åŸå› ": "", "ç¼ºæ¼": ", ".join(case["missing_fields"]),
                    "Claimåœ–ç‹€æ…‹": "N/A", "Claimåœ–èªªæ˜": ""
                }
                
                matched_pdf = None
                norm_case_key = normalize_string(case_key)
                
                # === V17 ä¿®æ­£ï¼šå¯¬é¬†æ¯”å°é‚è¼¯ (Smart Matching) ===
                for pdf_name, pdf_bytes in pdf_file_map.items():
                    norm_pdf_name = normalize_string(pdf_name)
                    
                    # 1. ç²¾æº–æ¯”å° (èˆŠé‚è¼¯)
                    if norm_case_key and ((norm_case_key in norm_pdf_name) or (norm_pdf_name in norm_case_key)):
                        if len(norm_case_key) > 5:
                            matched_pdf = pdf_bytes
                            break
                    
                    # 2. æ ¸å¿ƒæ•¸å­—æ¯”å° (æ–°é‚è¼¯)
                    # æå– Case ä¸­çš„ç´”æ•¸å­—: US11226533B2 -> 11226533
                    case_digits = re.sub(r'\D', '', case_key)
                    if len(case_digits) >= 4 and case_digits in norm_pdf_name:
                        matched_pdf = pdf_bytes
                        break
                
                if matched_pdf:
                    # 1. æŠ“å–ä¸»è¦ä»£è¡¨åœ–
                    img_list_main, msg_main = extract_images_from_pdf_v13(matched_pdf, target_fig, case_key, debug=debug_mode, log_prefix="[Main] ")
                    
                    if img_list_main:
                        case["image_list"] = img_list_main
                        status["ç‹€æ…‹"] = f"âœ… æˆåŠŸ ({len(img_list_main)}å¼µ)"
                        match_count += 1
                    else:
                        status["ç‹€æ…‹"] = "âš ï¸ ç¼ºåœ–"; status["åŸå› "] = msg_main

                    # 2. æŠ“å– Claim é™„åœ–
                    if add_claim_slide:
                        specific_claim_fig = parse_fig_number_from_claim(claim_text_content)
                        img_list_claim = []
                        msg_claim = ""
                        
                        if specific_claim_fig:
                            img_list_claim, msg_claim = extract_images_from_pdf_v13(matched_pdf, specific_claim_fig, case_key, debug=debug_mode, log_prefix="[Claim] ")
                            if img_list_claim:
                                status["Claimåœ–ç‹€æ…‹"] = f"âœ… å°ˆå±¬ ({len(img_list_claim)}å¼µ)"
                                status["Claimåœ–èªªæ˜"] = f"æ‰¾åˆ°æŒ‡å®šåœ–: {specific_claim_fig}"
                            else:
                                if img_list_main:
                                    img_list_claim = img_list_main
                                    status["Claimåœ–ç‹€æ…‹"] = "âš ï¸ æ²¿ç”¨ä¸»åœ–"
                                    status["Claimåœ–èªªæ˜"] = f"æŒ‡å®šåœ– ({specific_claim_fig}) æŠ“å–å¤±æ•—: {msg_claim}"
                                else:
                                    status["Claimåœ–ç‹€æ…‹"] = "âŒ ç¼ºåœ–"
                                    status["Claimåœ–èªªæ˜"] = "æŒ‡å®šå¤±æ•—ä¸”ç„¡ä¸»åœ–"
                        else:
                            if img_list_main:
                                img_list_claim = img_list_main
                                status["Claimåœ–ç‹€æ…‹"] = "âœ… åŒä¸»åœ–"
                                status["Claimåœ–èªªæ˜"] = "æœªæŒ‡å®š"
                            else:
                                status["Claimåœ–ç‹€æ…‹"] = "âŒ ç¼ºåœ–"
                                status["Claimåœ–èªªæ˜"] = "æœªæŒ‡å®šä¸”ç„¡ä¸»åœ–"
                        
                        case["claim_image_list"] = img_list_claim

                else:
                    if not target_fig: status["ç‹€æ…‹"] = "âš ï¸ ç¼ºè³‡è¨Š"; status["åŸå› "] = "Wordç„¡ä»£è¡¨åœ–"
                    else: status["ç‹€æ…‹"] = "âŒ ç„¡PDF"; status["åŸå› "] = f"æ‰¾ä¸åˆ°PDF: {case_key} (å·²å˜—è©¦å¯¬é¬†æ¯”å°)"
                status_report_list.append(status)

        if all_cases:
            st.session_state['slides_data'] = all_cases
            st.session_state['status_report'] = status_report_list
            st.success(f"å®Œæˆï¼å…± {len(all_cases)} ç­†è³‡æ–™ã€‚")
        else:
            st.warning("ç„¡è³‡æ–™ã€‚")

    if st.session_state['slides_data']:
        st.divider()
        if st.button("ğŸ—‘ï¸ æ¸…é™¤é‡ä¾†"):
            st.session_state['slides_data'] = []
            st.session_state['status_report'] = []
            st.session_state['debug_logs_map'] = {}
            st.rerun()

# --- ä¸»ç•«é¢ ---
if not st.session_state['slides_data']:
    st.info("ğŸ‘ˆ è«‹å…ˆä¸Šå‚³æª”æ¡ˆã€‚")
else:
    st.subheader(f"ğŸ“‹ é è¦½ (å·²æ’åº: ç”³è«‹äºº -> æ—¥æœŸ)")
    cols = st.columns(3)
    for i, data in enumerate(st.session_state['slides_data']):
        with cols[i % 3]:
            with st.container(border=True):
                st.markdown(f"**Case {i+1}**")
                st.caption(f"{data['clean_company']} | {data['clean_date']}")
                st.text(f"{data['clean_number']}")
                if data['image_list']:
                    st.image(data['image_list'][0], caption=f"ä¸»åœ– ({len(data['image_list'])})", use_column_width=True)
                
                if data.get('claim_image_list'):
                     st.image(data['claim_image_list'][0], caption=f"Claim ç”¨åœ– ({len(data['claim_image_list'])})", use_column_width=True)
                
                full_claim_text = data['claim_text']
                claims_preview = split_claims_text(full_claim_text)
                count_claims = len(claims_preview) if full_claim_text else 0
                st.caption(f"Claim: {count_claims} çµ„")

    def generate_ppt(slides_data, need_claim_slide):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        for data in slides_data:
            # === Main Slide ===
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            left, top, width, height = Inches(0.5), Inches(0.5), Inches(5.0), Inches(2.0)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame; tf.word_wrap = True
            p1 = tf.add_paragraph(); p1.text = f"å…¬é–‹è™Ÿï¼š{data['clean_number']}"; p1.font.size = Pt(20); p1.font.bold = True
            p2 = tf.add_paragraph(); p2.text = f"æ—¥æœŸï¼š{data['clean_date']}"; p2.font.size = Pt(20); p2.font.bold = True
            p3 = tf.add_paragraph(); p3.text = f"å…¬å¸ï¼š{data['clean_company']}"; p3.font.size = Pt(20); p3.font.bold = True

            img_left = Inches(5.5); img_top = Inches(0.5); img_width = Inches(7.0)
            img_list = data.get('image_list', [])
            
            if img_list:
                num_imgs = len(img_list)
                img_w = (7.0 / num_imgs) - 0.1
                img_h = 3.0
                for idx, img_bytes in enumerate(img_list):
                    this_left = 5.5 + (idx * (img_w + 0.1))
                    slide.shapes.add_picture(BytesIO(img_bytes), Inches(this_left), Inches(0.5), height=Inches(img_h))
                
                text_top = Inches(3.6)
                text_height = Inches(1.0)
                txBox = slide.shapes.add_textbox(img_left, text_top, img_width, text_height)
                tf = txBox.text_frame; tf.word_wrap = True
                content = data['rep_fig_text'] if data['rep_fig_text'].strip() else ""
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph(); p.text = line.strip(); p.font.size = Pt(14)
            else:
                img_height = Inches(4.0)
                txBox = slide.shapes.add_textbox(img_left, img_top, img_width, img_height)
                tf = txBox.text_frame; tf.word_wrap = True
                content = data['rep_fig_text'] if data['rep_fig_text'].strip() else "ç„¡ä»£è¡¨åœ–è³‡è¨Š"
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph(); p.text = line.strip(); p.font.size = Pt(16)

            left, top, width, height = Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame; tf.word_wrap = True
            p1 = tf.add_paragraph(); p1.text = "â€¢ è§£æ±ºå•é¡Œï¼š" + data['problem']; p1.font.size = Pt(18); p1.space_after = Pt(12)
            p2 = tf.add_paragraph(); p2.text = "â€¢ ç™¼æ˜ç²¾ç¥ï¼š" + data['spirit']; p2.font.size = Pt(18)

            left, top, width, height = Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 192, 0); shape.line.color.rgb = RGBColor(255, 192, 0)
            p = shape.text_frame.paragraphs[0]; p.text = data['key_point']; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(20); p.font.bold = True
            shape.text_frame.vertical_anchor = MSO_SHAPE.RECTANGLE

            # === Claim Slides ===
            if need_claim_slide:
                claims_groups = split_claims_text(data['claim_text'])
                if not claims_groups and data['claim_text'].strip():
                      claims_groups = [data['claim_text'].split('\n')]

                for claim_lines in claims_groups:
                    slide_c = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    left, top, width, height = Inches(0.5), Inches(0.5), Inches(5.0), Inches(2.0)
                    txBox = slide_c.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame; tf.word_wrap = True
                    p1 = tf.add_paragraph(); p1.text = f"å…¬é–‹è™Ÿï¼š{data['clean_number']}"; p1.font.size = Pt(20); p1.font.bold = True
                    p2 = tf.add_paragraph(); p2.text = f"æ—¥æœŸï¼š{data['clean_date']}"; p2.font.size = Pt(20); p2.font.bold = True
                    p3 = tf.add_paragraph(); p3.text = f"å…¬å¸ï¼š{data['clean_company']}"; p3.font.size = Pt(20); p3.font.bold = True
                    
                    # è²¼ä¸Š Claim åœ–ç‰‡ (å¦‚æœæœ‰)
                    claim_imgs = data.get('claim_image_list', [])
                    if claim_imgs:
                        img_left = Inches(5.5); img_top = Inches(0.5)
                        num_imgs = len(claim_imgs)
                        img_w = (7.0 / num_imgs) - 0.1
                        img_h = 3.0
                        for idx, img_bytes in enumerate(claim_imgs):
                            this_left = 5.5 + (idx * (img_w + 0.1))
                            slide_c.shapes.add_picture(BytesIO(img_bytes), Inches(this_left), Inches(0.5), height=Inches(img_h))

                    # æ ¹æ“šæœ‰ç„¡åœ–ç‰‡èª¿æ•´æ–‡å­—æ¡†ä½ç½®
                    left, width = Inches(0.5), Inches(12.3)
                    if claim_imgs:
                         top = Inches(3.6); height = Inches(3.4)
                    else:
                         top = Inches(2.5); height = Inches(4.5)

                    txBox = slide_c.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame; tf.word_wrap = True
                    
                    p_title = tf.add_paragraph()
                    p_title.text = "ã€ç¨ç«‹é … Claimã€‘"
                    p_title.font.size = Pt(24); p_title.font.bold = True; p_title.font.color.rgb = RGBColor(0, 112, 192)
                    p_title.space_after = Pt(10)
                    
                    for line in claim_lines:
                        clean_line = line.strip()
                        if clean_line:
                            p = tf.add_paragraph()
                            p.text = clean_line
                            p.font.size = Pt(14) 
                            p.space_after = Pt(4)
                            
                            if line.startswith('\t') or line.startswith('    '):
                                p.level = 1
                            elif clean_line.startswith(('o ', 'â—‹', '-', 'â€¢', 'â—')):
                                p.level = 1
                            elif clean_line.startswith(('â–ª', 'â– ')):
                                p.level = 2
                            elif re.match(r'^(\(\d+\)|\d+\.|\d+\))', clean_line):
                                if "Claim" in clean_line or "ç¨ç«‹é …" in clean_line:
                                    p.level = 0
                                    p.font.bold = True
                                else:
                                    p.level = 1

        return prs

    st.divider()
    if st.button("ğŸš€ ç”Ÿæˆ PowerPoint (.pptx)", type="primary"):
        prs = generate_ppt(st.session_state['slides_data'], add_claim_slide)
        binary_output = BytesIO()
        prs.save(binary_output)
        binary_output.seek(0)
        st.download_button("ğŸ“¥ ä¸‹è¼‰ PPT", binary_output, "slides_with_claims.pptx")

    st.divider()
    st.subheader("ğŸ“Š è¨ºæ–·å ±å‘Š")
    if st.session_state['status_report']:
        df = pd.DataFrame(st.session_state['status_report'])
        cols = ["ä¾†æº", "æ¡ˆè™Ÿ(å…¬é–‹è™Ÿ)", "å…¬å¸", "æ—¥æœŸ(å„ªå…ˆæ¬Šæ—¥)", "å°æ‡‰PPTçš„é ç¢¼", "ç‹€æ…‹", "åŸå› ", "Claimåœ–ç‹€æ…‹", "Claimåœ–èªªæ˜", "ç¼ºæ¼"]
        st.dataframe(df[cols], hide_index=True)
